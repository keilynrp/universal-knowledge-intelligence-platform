"""
Report export format parity — validation consistency (change:
unify-report-format-coverage, phase 5).

Two standalone defects, independent of the data/presentation refactor:
  1. /exports/excel is the only export endpoint that never validated section
     names, so a typo or stale client is accepted silently.
  2. The PPTX exporter gates its top-labels slide on the deprecated `top_brands`
     alias, which GET /reports/sections does not return — a client using the
     documented public id `top_secondary_labels` silently loses the slide.
"""
import io

from pptx import Presentation

from backend import models


_UNKNOWN = "definitely_not_a_section"


# ── 5.1 / 5.3 — Excel validates unknown sections, like the other three ─────────

def test_excel_rejects_unknown_section(client, auth_headers):
    """POST /exports/excel with an unknown section must 422, not silently 200."""
    payload = {"domain_id": "default", "sections": ["entity_stats", _UNKNOWN]}
    resp = client.post("/exports/excel", json=payload, headers=auth_headers)
    assert resp.status_code == 422
    assert _UNKNOWN in resp.text


def test_all_export_endpoints_reject_unknown_section_identically(client, auth_headers):
    """HTML, PDF, Excel and PPTX must all reject the same unknown name with 422."""
    payload = {"domain_id": "default", "sections": [_UNKNOWN]}
    for path in ("/reports/generate", "/exports/pdf", "/exports/excel", "/exports/pptx"):
        resp = client.post(path, json=payload, headers=auth_headers)
        assert resp.status_code == 422, f"{path} accepted an unknown section"


# ── 5.4 — PPTX honours the public section id, not only the deprecated alias ────

def test_pptx_renders_top_labels_for_public_id(client, auth_headers, db_session):
    """The top-labels slide must render for the PUBLIC id `top_secondary_labels`.

    RED: pptx_exporter gates on `top_brands`, the deprecated alias that
    /reports/sections does not return, so a documented-vocabulary client gets
    no slide.
    """
    db_session.add(models.RawEntity(
        primary_label="Labelled record",
        domain="default",
        secondary_label="Clinical Trial",
    ))
    db_session.commit()

    payload = {"domain_id": "default", "sections": ["top_secondary_labels"]}
    resp = client.post("/exports/pptx", json=payload, headers=auth_headers)
    assert resp.status_code == 200

    prs = Presentation(io.BytesIO(resp.content))
    texts = [
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]
    blob = "\n".join(texts)
    assert "Top Secondary Labels" in blob, "top-labels slide missing for public id"
    assert "Clinical Trial" in blob
