"""Per-format renderers over the section payload (unify-report-format-coverage,
phase 2).

Each renderer turns a format-neutral `SectionData` into one output format.
These tests exercise every block type against each renderer in isolation, with
no database and no migration wiring — they pin that the primitives render
faithfully before any section is migrated onto them.
"""
import io

from pptx import Presentation

from backend.reporting.section_data import (
    Meter,
    Narrative,
    SectionData,
    StatGrid,
    StatItem,
    Table,
)
import pytest

pytestmark = pytest.mark.reporting


def _every_block_section() -> SectionData:
    return SectionData(
        key="demo",
        title="Demo Section",
        takeaway="1,240 entities recorded, 60% enriched",
        method="Counts scoped to this domain, as of the last run.",
        blocks=(
            StatGrid(items=(
                StatItem(label="Total", value="1,240"),
                StatItem(label="Enriched", value="60%", sub="744 of 1,240"),
            )),
            Table(
                columns=("Status", "Count", "Share"),
                rows=(("Validated", "800", "64"), ("Pending", "440", "36")),
                bar_column=2,
            ),
            Narrative(heading="Executive reading", paragraphs=("First point.", "Second point.")),
            Meter(label="Coverage", pct=60),
        ),
    )


# ── HTML renderer (2.1) ───────────────────────────────────────────────────────

def test_html_renderer_emits_existing_css_classes():
    from backend.reporting.html_renderer import render_html

    html = render_html(_every_block_section())

    # Section wrapper + title
    assert "<section>" in html
    assert "Demo Section" in html
    # StatGrid → grid / stat-card / label / value / sub
    for cls in ("class=\"grid\"", "class=\"stat-card\"", "class=\"label\"", "class=\"value\"", "class=\"sub\""):
        assert cls in html, f"missing {cls}"
    assert "744 of 1,240" in html
    # Table → table / th / cells
    assert "<table>" in html and "<th>" in html
    assert "Validated" in html and "800" in html
    # Table bar_column → bar markup for that column
    assert "class=\"bar-wrap\"" in html and "class=\"bar\"" in html
    # Narrative → callout / h3 / paragraphs
    assert "class=\"callout\"" in html and "<h3>" in html
    assert "Executive reading" in html and "First point." in html
    # Meter → bar markup with pct
    assert "Coverage" in html and "60%" in html


def test_html_renderer_leads_with_the_takeaway_not_the_label():
    """The heading states the finding; the dataset label drops to secondary text.

    Both are present — a report a reader cannot scan by section name is worse
    than one that only names its sections — but the assertion is what leads.
    """
    from backend.reporting.html_renderer import render_html

    html = render_html(_every_block_section())

    assert "<h2>1,240 entities recorded, 60% enriched</h2>" in html
    assert "<h2>Demo Section</h2>" not in html, "the label is still the heading"
    assert "Demo Section" in html, "the label has to remain findable"


def test_html_renderer_states_the_method_for_every_section():
    from backend.reporting.html_renderer import render_html

    html = render_html(_every_block_section())

    assert 'class="method"' in html
    assert "Counts scoped to this domain, as of the last run." in html


def test_html_renderer_shows_the_exhibit_ordinal_assembly_assigned():
    from dataclasses import replace

    from backend.reporting.html_renderer import render_html

    html = render_html(replace(_every_block_section(), exhibit=3))

    assert "Exhibit 3" in html


def test_html_renderer_omits_the_ordinal_for_an_unnumbered_section():
    """A section rendered outside assembly has no ordinal, and must not invent
    one or leak the absent value."""
    from backend.reporting.html_renderer import render_html

    html = render_html(_every_block_section())  # exhibit is None

    assert "Exhibit" not in html
    assert "None" not in html


def test_html_renderer_escapes_the_takeaway_and_method():
    """Both are prose assembled from record values, so both can carry markup."""
    from backend.reporting.html_renderer import render_html

    html = render_html(SectionData(
        key="x", title="T",
        takeaway="<script>alert(1)</script> & up",
        method="Source: <b>upstream</b> & local",
    ))

    assert "<script>" not in html and "<b>" not in html
    assert "&lt;script&gt;" in html and "&amp;" in html


def test_html_renderer_escapes_data():
    from backend.reporting.html_renderer import render_html

    section = SectionData(
        key="x", title="T",
        takeaway="a finding", method="a source",
        blocks=(StatGrid(items=(StatItem(label="<script>", value="a&b"),)),),
    )
    html = render_html(section)
    assert "<script>" not in html
    assert "&lt;script&gt;" in html
    assert "a&amp;b" in html


# ── Excel renderer (2.3) ──────────────────────────────────────────────────────

def _excel_cell_blob(ws) -> str:
    return "\n".join(
        str(cell.value) for row in ws.iter_rows() for cell in row if cell.value is not None
    )


def test_excel_renderer_writes_all_block_types():
    import openpyxl
    from backend.reporting.excel_renderer import render_excel

    wb = openpyxl.Workbook()
    ws = render_excel(_every_block_section(), wb)

    # A sheet named for the section (Excel truncates to 31 chars).
    assert ws.title == "Demo Section"
    assert ws.title in wb.sheetnames

    blob = _excel_cell_blob(ws)
    # StatGrid → KPI block
    assert "Total" in blob and "1,240" in blob and "744 of 1,240" in blob
    # Table → columns + rows
    assert "Status" in blob and "Validated" in blob and "800" in blob
    # Narrative → heading + paragraph
    assert "Executive reading" in blob and "First point." in blob
    # Meter → percentage cell
    assert "Coverage" in blob and "60%" in blob


def _cell_column(ws, col: int = 1) -> list[str]:
    return [str(c.value) for c in ws[chr(ord("A") + col - 1)] if c.value is not None]


def test_excel_sheet_opens_with_the_finding():
    import openpyxl
    from backend.reporting.excel_renderer import render_excel

    ws = render_excel(_every_block_section(), openpyxl.Workbook())

    assert ws["A1"].value == "1,240 entities recorded, 60% enriched"
    assert ws["A1"].font.bold


def test_excel_carries_the_caveat_directly_above_each_table():
    """5.3 — Excel is the format most often re-cut and pasted, and that is the
    path by which a proxy metric loses its caveat. The disclosure sits on the
    row immediately above the header, so it travels with the range."""
    import openpyxl
    from backend.reporting.excel_renderer import render_excel

    ws = render_excel(_every_block_section(), openpyxl.Workbook())

    header_row = next(
        cell.row
        for row in ws.iter_rows()
        for cell in row
        if cell.value == "Status"
    )
    above = ws.cell(row=header_row - 1, column=1).value
    assert above and "Counts scoped to this domain" in above, above


def test_excel_discloses_a_section_that_has_no_table():
    """Nothing to attach the caveat to, so it goes under the finding — the
    disclosure is mandatory, not conditional on the block types present."""
    import openpyxl
    from backend.reporting.excel_renderer import render_excel

    section = SectionData(
        key="k", title="Narrative Only",
        takeaway="Three patterns detected",
        method="Statistical co-occurrence within this corpus, not causation.",
        blocks=(Narrative(heading="Reading", paragraphs=("A point.",)),),
    )
    ws = render_excel(section, openpyxl.Workbook())

    blob = "\n".join(_cell_column(ws))
    assert "not causation" in blob


def test_excel_renderer_truncates_long_sheet_names():
    import openpyxl
    from backend.reporting.excel_renderer import render_excel

    long_title = "A very long section title that exceeds the Excel limit"
    wb = openpyxl.Workbook()
    ws = render_excel(SectionData(key="k", title=long_title, takeaway="t", method="m", blocks=()), wb)
    assert len(ws.title) <= 31


def test_excel_renderer_sanitizes_invalid_sheet_chars():
    import openpyxl
    from backend.reporting.excel_renderer import render_excel

    wb = openpyxl.Workbook()
    ws = render_excel(SectionData(key="k", title="A/B:C*?[D]", takeaway="t", method="m", blocks=()), wb)
    assert not (set(ws.title) & set(r"[]:*?/\\"))


# ── PPTX renderer (2.5) ───────────────────────────────────────────────────────

def _pptx_text(prs) -> str:
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
            if shape.has_table:
                for r in shape.table.rows:
                    for c in r.cells:
                        out.append(c.text)
    return "\n".join(out)


def _has_table(prs) -> bool:
    return any(shape.has_table for slide in prs.slides for shape in slide.shapes)


def test_pptx_renderer_renders_all_block_types():
    from pptx import Presentation as _P
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from backend.reporting.pptx_renderer import render_pptx

    prs = _P()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slides = render_pptx(_every_block_section(), prs, RGBColor(0x63, 0x66, 0xF1))

    assert slides, "renderer added no slides"
    blob = _pptx_text(prs)
    # Section title
    assert "Demo Section" in blob
    # StatGrid → KPI row
    assert "Total" in blob and "1,240" in blob
    # Table → a real pptx table with headers + cells
    assert _has_table(prs)
    assert "Status" in blob and "Validated" in blob
    # Narrative → bullet slide
    assert "Executive reading" in blob and "First point." in blob
    # Meter → label + pct
    assert "Coverage" in blob and "60%" in blob


def test_pptx_renderer_draws_a_bar_shape_for_meter():
    from pptx import Presentation as _P
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from backend.reporting.pptx_renderer import render_pptx

    prs = _P()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    section = SectionData(key="k", title="Only Meter", takeaway="t", method="m", blocks=(Meter(label="Coverage", pct=42),))
    render_pptx(section, prs, RGBColor(0x63, 0x66, 0xF1))

    auto_shapes = [
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    ]
    assert auto_shapes, "meter produced no bar shape"


# ── Cross-cutting: block support + round-trip (2.7, 1.3) ──────────────────────

def test_every_renderer_declares_support_for_all_four_blocks():
    from backend.reporting import html_renderer, excel_renderer, pptx_renderer

    expected = {StatGrid, Table, Narrative, Meter}
    for module in (html_renderer, excel_renderer, pptx_renderer):
        assert module.SUPPORTED_BLOCKS == expected, module.__name__


def test_every_block_type_round_trips_through_each_renderer():
    """A payload using every block type renders in all three formats, no raises."""
    import openpyxl
    from pptx import Presentation as _P
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from backend.reporting.html_renderer import render_html
    from backend.reporting.excel_renderer import render_excel
    from backend.reporting.pptx_renderer import render_pptx

    section = _every_block_section()

    html = render_html(section)
    assert html.startswith("<section>")

    ws = render_excel(section, openpyxl.Workbook())
    assert ws.max_row >= 1

    prs = _P()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slides = render_pptx(section, prs, RGBColor(0x63, 0x66, 0xF1))
    assert slides
