"""Parity guard — the definition of done for unify-report-format-coverage.

Every public section must either render in every export format or be explicitly
declared unsupported *and reported as omitted* — never silently dropped. The
migration is complete: there are no xfails. Each (format, section) combo is now
asserted directly against `SECTION_FORMAT_SUPPORT`:

  * A combo the support map claims → the body performs the *real* render and the
    section's marker must appear, or the test fails (the map cannot lie about
    coverage it does not have).
  * A combo the map omits → the export must still succeed (no raise) and the
    section must be reported by `unsupported_sections()` — the omission contract
    (`test_report_omissions`), not a silent drop.

The map and reality cannot drift apart in either direction: claiming support
without a renderer fails the render assertion; rendering without claiming it
means the "unsupported" branch finds the marker anyway is impossible because the
branch is only taken when the map omits the combo. agentic_trace is the sole
declared-unsupported section (Excel + PPTX; free text that belongs in HTML/PDF).
"""
import io
import re
from html import escape, unescape

import pytest
from pptx import Presentation

from backend import models, report_builder
from backend.reporting.localize import localize_section
from backend.exporters.excel_exporter import EnterpriseExcelExporter
from backend.exporters.pptx_exporter import generate_pptx
from backend.reporting import format_support
from backend.reporting.section_data import Table as PayloadTable

pytestmark = pytest.mark.reporting


# How each section is expected to appear, per format. For HTML/PDF it is the
# section's <h2> label; for Excel it is the sheet name; for PPTX the slide
# title. A section absent from a format's map has no representation there yet —
# adding a renderer for it (migration) means adding its marker here too.
_MARKERS: dict[str, dict[str, str]] = {
    "html": dict(report_builder.SECTION_LABELS),
    "pdf": dict(report_builder.SECTION_LABELS),
    "excel": {
        "entity_stats": "Entity Statistics",  # migrated (phase 3): dedicated sheet
        "enrichment_coverage": "Enrichment Coverage",  # migrated (phase 3.2)
        "top_secondary_labels": "Top Secondary Labels",  # migrated (phase 3.3)
        "impact_projection": "Impact Projection",  # migrated (phase 3.7)
        "institutional_benchmark": "Institutional Benchmark",  # migrated (phase 3.6)
        "hidden_patterns": "Hidden Patterns",  # migrated (phase 3.8)
        "decision_recommendations": "Suggested Next Actions",  # migrated (phase 3.9)
        "topic_clusters": "Top Concepts",  # migrated 3.4, relabelled 3.6
        "harmonization_log": "Harmonization",
        "authority_control": "Authority Control",  # extend-report-module-coverage
        "collaboration_graph": "Collaboration Graph",
        "journal_portfolio": "Journal Portfolio",
    },
    "pptx": {
        "entity_stats": "Entity Statistics",
        "enrichment_coverage": "Enrichment Coverage",
        "top_secondary_labels": "Top Secondary Labels",
        "topic_clusters": "Top Concepts",  # migrated 3.4, relabelled 3.6
        # migrated (phase 3): rendered via the shared payload + render_pptx
        "impact_projection": "Impact Projection",
        "institutional_benchmark": "Institutional Benchmark",
        "hidden_patterns": "Hidden Patterns",
        "decision_recommendations": "Suggested Next Actions",
        "harmonization_log": "Harmonization Log",
        "authority_control": "Authority Control",  # extend-report-module-coverage
        "collaboration_graph": "Collaboration Graph",
        "journal_portfolio": "Journal Portfolio",
    },
}

_BRANDING = {
    "platform_name": "UKIP",
    "logo_url": None,
    "accent_color": "#6366f1",
    "footer_text": "UKIP",
}


def _render(export_format: str, section: str, db) -> str:
    """Render a single section in one format, returned as a searchable blob.

    Every blob carries everything a reader of that format can see, not just the
    part that identifies the section: sheet *names and cells* for Excel, text
    frames *and speaker notes* for PPTX. The presentation check (6.2) needs the
    contents, and a marker check over a blob that omits them would be one more
    gate that cannot fail.
    """
    if export_format in ("html", "pdf"):
        # PDF renders exactly this HTML through WeasyPrint, so the HTML blob is
        # the faithful check for both without invoking a native PDF engine.
        return report_builder.build(db, "default", [section])
    if export_format == "excel":
        data = EnterpriseExcelExporter().build(db, "default", [section])
        wb = load_workbook_from_bytes(data)
        cells = [
            str(cell.value)
            for ws in wb.worksheets
            for row in ws.iter_rows()
            for cell in row
            if cell.value is not None
        ]
        return "\n".join([*wb.sheetnames, *cells])
    if export_format == "pptx":
        data = generate_pptx(
            db=db, domain_id="default", sections=[section], title=None,
            branding=_BRANDING, org_id=None,
        )
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            parts += [
                shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
            ]
            if slide.has_notes_slide:
                parts.append(slide.notes_slide.notes_text_frame.text)
        return "\n".join(parts)
    raise AssertionError(f"unknown format {export_format}")


def _visible(export_format: str, blob: str, text: str) -> bool:
    """Whether `text` appears anywhere in the format's output.

    Deliberately weak, and used only for the negative case in 6.3. Presence in
    the file is NOT the contract and must not be asserted as if it were: mutation
    testing found that searching the whole output passes even when a section
    states nothing, because the executive summary already carries every takeaway,
    the `Methodology` sheet carries every method, and PPTX speaker notes carry
    both. Four of six deliberate breakages went undetected that way.

    Placement is the contract. See `_placement`.
    """
    if export_format in ("html", "pdf"):
        return escape(text) in blob
    return text in blob


#: Sheets the Excel exporter always writes, whatever sections were requested.
_EXCEL_INFRA_SHEETS = frozenset({"Summary", "Methodology", "Entities"})


def _placement(export_format: str, section: str, db, payload) -> dict[str, object]:
    """What the format actually put in each slot the contract names for it.

    One extractor per format, reading the section's *own* rendering rather than
    the whole document: the `<section>` element for HTML, the section's worksheet
    for Excel, the section's slides for PPTX. Returns None for a slot the format
    left empty, so the caller can say which slot is missing rather than just
    "the text is not in the file".
    """
    label = report_builder.SECTION_LABELS[section]

    if export_format in ("html", "pdf"):
        doc = report_builder.build(db, "default", [section])
        for block in re.findall(r"<section>.*?</section>", doc, re.S):
            eyebrow = re.search(r'<div class="exhibit-label">(.*?)</div>', block, re.S)
            if not eyebrow or escape(label) not in eyebrow.group(1):
                continue
            heading = re.search(r"<h2>(.*?)</h2>", block, re.S)
            method = re.search(r'<p class="method">(.*?)</p>', block, re.S)
            return {
                "takeaway": unescape(heading.group(1)) if heading else None,
                "method": unescape(method.group(1)) if method else None,
                "exhibit": unescape(eyebrow.group(1)),
            }
        return {"takeaway": None, "method": None, "exhibit": None}

    if export_format == "excel":
        wb = load_workbook_from_bytes(
            EnterpriseExcelExporter().build(db, "default", [section])
        )
        names = [n for n in wb.sheetnames if n not in _EXCEL_INFRA_SHEETS]
        assert len(names) == 1, f"expected one section sheet, got {names}"
        ws = wb[names[0]]

        # Find the caveat, then check what it sits above — rather than finding
        # the table and checking what sits above it. The payload's column names
        # cannot locate the header on the one sheet still written by hand, whose
        # columns differ from its payload's, and adjacency is the requirement
        # either way: a copied range has to carry its warning.
        def row_width(row: int) -> int:
            return sum(1 for cell in ws[row] if cell.value is not None)

        method_rows = [
            cell.row for cell in ws["A"] if str(cell.value) == payload.method
        ]
        return {
            "takeaway": ws["A1"].value,
            "method": payload.method if method_rows else None,
            # A caveat is adjacent when the row under it is a header — two or more
            # cells wide. Used only to *locate*; whether adjacency is required at
            # all is decided from the payload, which knows definitively. Scanning
            # the sheet for "a table" instead reads a StatGrid row (label, value,
            # sub) as one, and then demands adjacency of sections that have no
            # table to be adjacent to.
            "adjacent": any(row_width(r + 1) >= 2 for r in method_rows),
            "sheet_has_table": any(
                isinstance(block, PayloadTable) for block in payload.blocks
            ),
            "exhibit": "\n".join(
                str(c.value) for r in ws.iter_rows() for c in r if c.value is not None
            ),
        }

    if export_format == "pptx":
        prs = Presentation(io.BytesIO(generate_pptx(
            db=db, domain_id="default", sections=[section], title=None,
            branding=_BRANDING, org_id=None,
        )))
        frames, notes = [], []
        for slide in prs.slides:
            texts = [
                shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
            ]
            if not any(label in t for t in texts):
                continue  # cover / closing slide
            frames += texts
            if slide.has_notes_slide:
                notes.append(slide.notes_slide.notes_text_frame.text)
        # The title is set verbatim, so an exact frame match is the assertion —
        # a slide titled with the label instead would have no such frame.
        title = next((t for t in frames if t == payload.takeaway), None)
        # The footer is clipped by design, so it is checked as a prefix with the
        # ellipsis removed rather than against a hard-coded character limit.
        footer = next(
            (
                t for t in frames
                if t and payload.method.startswith(t.removesuffix(" …").rstrip())
                and t != payload.takeaway
            ),
            None,
        )
        return {
            "takeaway": title,
            "method": footer,
            "notes": "\n".join(notes),
            "exhibit": "\n".join(frames),
        }

    raise AssertionError(f"unknown format {export_format}")


def load_workbook_from_bytes(data: bytes):
    import openpyxl
    return openpyxl.load_workbook(io.BytesIO(data))


def _seed(db) -> None:
    """Enough data for every section to produce its real content."""
    for idx in range(3):
        db.add(models.RawEntity(
            primary_label=f"Parity record {idx}",
            domain="default",
            enrichment_status="completed",
            enrichment_concepts="knowledge graph; semantic intelligence",
            enrichment_citation_count=120 + idx,
            enrichment_source="openalex",
            secondary_label="Clinical Trial",
            quality_score=0.8,
        ))
    db.add(models.HarmonizationLog(
        step_id="normalize_labels",
        step_name="Normalize labels",
        records_updated=3,
        fields_modified="primary_label",
    ))
    # Authority records so authority_control renders its real content (KPI grid,
    # distribution, conflicts) rather than its empty state — a section that only
    # ever renders "not available" would pass the marker check without proving
    # the populated path works in every format.
    db.add(models.AuthorityRecord(
        field_name="brand_capitalized", original_value="acme corp",
        canonical_label="ACME Corporation", confidence=0.92,
        status="confirmed", resolution_status="exact_match", review_required=False,
    ))
    db.add(models.AuthorityRecord(
        field_name="brand_capitalized", original_value="initech",
        confidence=0.41, status="pending", resolution_status="ambiguous",
        review_required=True, nil_reason="multiple_candidates",
    ))
    # Author stats + a cross-community edge so collaboration_graph renders its
    # populated content (counts, centrality, a bridge) in every format.
    for aid, key, name, comm, cent in [
        (1, "alice", "Alice Ng", 1, 0.9), (2, "bob", "Bob Ito", 1, 0.5),
        (3, "carol", "Carol Vex", 2, 0.7),
    ]:
        db.add(models.Author(id=aid, name_key=key, display_name=name))
        db.add(models.AuthorStats(
            author_id=aid, org_id=None, domain_id="default",
            degree=3, centrality=cent, community_id=comm, publication_count=10,
        ))
    db.add(models.CoauthorEdge(author_a_id=2, author_b_id=3, org_id=None,
                               domain_id="default", weight=1.0))
    # A journal metric with a full credible interval so journal_portfolio renders
    # its populated table (NIF, Bayesian NIF [CI], DOAJ, APC) in every format.
    db.add(models.JournalMetric(
        org_id=None, issn_l="issn-parity", display_name="Parity Journal",
        normalized_impact_factor=2.10, nif_field="cs",
        nif_bayes=2.05, nif_ci_low=1.60, nif_ci_high=2.50,
        works_2yr=8, apc_usd=1500, is_in_doaj=True,
    ))
    db.commit()


_COMBOS = [
    pytest.param(fmt, section, id=f"{fmt}:{section}")
    for fmt in format_support.EXPORT_FORMATS
    for section in format_support.PUBLIC_SECTIONS
]


@pytest.mark.parametrize("export_format,section", _COMBOS)
def test_section_renders_or_is_declared_unsupported(export_format, section, db_session):
    _seed(db_session)
    blob = _render(export_format, section, db_session)  # must never raise
    if format_support.supports(export_format, section):
        marker = _MARKERS[export_format].get(section)
        assert marker is not None, f"{export_format} declares no marker for {section}"
        assert marker in blob, f"{export_format} did not render {section} (marker {marker!r})"
    else:
        # Declared unsupported: the export still succeeds and the section is
        # reported as omitted (see test_report_omissions), never silently dropped.
        assert section in format_support.unsupported_sections(export_format, [section])


# ── Presentation parity (report-presentation 6.1–6.3) ────────────────────────


def test_every_format_declares_the_mandatory_presentation_elements():
    """6.1 — the declaration is itself constrained.

    Section coverage is a ratchet a format may sit below. Presentation coverage
    is not: takeaway and method are not declarable as unsupported, so a future
    format cannot opt out of the contract by leaving them out of its entry.
    """
    assert set(format_support.PRESENTATION_SUPPORT) == set(format_support.EXPORT_FORMATS), (
        "a format is missing from PRESENTATION_SUPPORT, or invented in it"
    )
    for fmt, elements in format_support.PRESENTATION_SUPPORT.items():
        missing = format_support.REQUIRED_PRESENTATION_ELEMENTS - elements
        assert not missing, f"{fmt} declares {sorted(missing)} unsupported"
        unknown = elements - set(format_support.PRESENTATION_ELEMENTS)
        assert not unknown, f"{fmt} claims unknown elements: {sorted(unknown)}"


@pytest.mark.parametrize("export_format,section", _COMBOS)
def test_a_rendered_section_states_its_finding_and_its_method(
    export_format, section, db_session
):
    """6.2 — presentation coverage cannot drift from section coverage.

    A format that renders a section must emit that section's takeaway and its
    method disclosure. Without this, every earlier phase is reversible by
    accident: a renderer refactor that drops the footer would leave the section
    still rendering, still passing the marker check, and silently uncited.

    Compared against the payload the collector produced rather than against
    hard-coded strings, so an edited takeaway does not need this test edited too —
    and `collect_section` is the same dispatch `build()` uses, not a second copy
    of it.
    """
    _seed(db_session)
    if not format_support.supports(export_format, section):
        pytest.skip(f"{export_format} declares {section} unsupported — see 6.3")

    # Localized before comparison: since #268 a collector emits catalog keys and
    # the renderer resolves them, so "the payload the collector produced" means
    # the payload after localization. Comparing against the raw one would ask the
    # renderer to emit a key. The guard's strength is unchanged — it still fails
    # if a renderer drops the method — and it still avoids hard-coded strings.
    payload = localize_section(
        report_builder.collect_section(db_session, section, "default")
    )
    assert payload is not None, f"{section} has no collector"

    slots = _placement(export_format, section, db_session, payload)

    assert slots["takeaway"] == payload.takeaway, (
        f"{export_format} does not lead {section} with its finding.\n"
        f"expected: {payload.takeaway!r}\nin slot:  {slots['takeaway']!r}"
    )
    assert slots["method"], (
        f"{export_format} renders {section} with an empty method slot"
    )
    if export_format == "pptx":
        # Footer clipped, notes complete — both, which is why clipping is safe.
        assert payload.method.startswith(str(slots["method"]).removesuffix(" …").rstrip())
        assert payload.method in slots["notes"], (
            "the speaker notes must carry the disclosure untruncated"
        )
    elif export_format == "excel":
        assert slots["method"] == payload.method
        if slots["sheet_has_table"]:
            assert slots["adjacent"], (
                f"{section}'s caveat is on the sheet but not directly above a "
                "table header, so a copied range would leave it behind"
            )
    else:
        assert slots["method"] == payload.method, (
            f"{export_format} does not disclose {section}'s method where the "
            f"contract puts it.\nexpected: {payload.method!r}\n"
            f"in slot:  {slots['method']!r}"
        )

    # Decision 7: the ordinal is a within-document reference and belongs only to
    # the formats that are documents.
    if format_support.carries(export_format, "exhibit"):
        assert "Exhibit" in str(slots["exhibit"])
    else:
        assert "Exhibit" not in str(slots["exhibit"]), (
            f"{export_format} does not declare `exhibit` but numbered {section}"
        )


@pytest.mark.parametrize(
    "export_format", [f for f in format_support.EXPORT_FORMATS if f not in ("html", "pdf")]
)
def test_an_unsupported_section_is_exempt_and_still_reported(export_format, db_session):
    """6.3 — the exemption is real, and it is an omission rather than a partial.

    `agentic_trace` is the only declared-unsupported section. The format must not
    be held to the presentation contract for it, must still be asked to render
    without raising, and must still name it as omitted — and "unsupported" has to
    mean the section is absent, not present-without-its-statements, which would
    be the worst of both.
    """
    _seed(db_session)
    section = "agentic_trace"
    assert not format_support.supports(export_format, section), (
        f"{section} is now supported by {export_format}; this test needs a new subject"
    )

    # Localized before comparison: since #268 a collector emits catalog keys and
    # the renderer resolves them, so "the payload the collector produced" means
    # the payload after localization. Comparing against the raw one would ask the
    # renderer to emit a key. The guard's strength is unchanged — it still fails
    # if a renderer drops the method — and it still avoids hard-coded strings.
    payload = localize_section(
        report_builder.collect_section(db_session, section, "default")
    )
    blob = _render(export_format, section, db_session)  # must not raise

    assert not _visible(export_format, blob, payload.takeaway), (
        f"{export_format} declares {section} unsupported but rendered its takeaway"
    )
    # The omission contract (test_report_omissions) is untouched by any of this.
    assert section in format_support.unsupported_sections(export_format, [section])


def test_support_map_covers_every_public_section_key():
    """Every format lists only real public sections; no typos, no aliases."""
    for fmt, sections in format_support.SECTION_FORMAT_SUPPORT.items():
        unknown = sections - set(format_support.PUBLIC_SECTIONS)
        assert not unknown, f"{fmt} claims unknown sections: {unknown}"


def test_builder_and_collector_maps_cannot_drift():
    """The two section maps must stay key-identical.

    SECTION_BUILDERS is no longer what build() assembles from — SECTION_COLLECTORS
    is — but it remains the section registry: four export endpoints validate
    requested names against it, scheduled reports filter on it, and
    format_support derives PUBLIC_SECTIONS from it.

    So a section added to one map and not the other would either be requestable
    and unrenderable, or renderable and rejected as unknown. This session has
    already paid for two maps drifting apart in exactly this way.
    """
    from backend import report_builder

    builders = set(report_builder.SECTION_BUILDERS)
    collectors = set(report_builder.SECTION_COLLECTORS)
    assert builders == collectors, (
        f"only in SECTION_BUILDERS: {sorted(builders - collectors)}; "
        f"only in SECTION_COLLECTORS: {sorted(collectors - builders)}"
    )
