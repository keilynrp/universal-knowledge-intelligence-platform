"""The presentation contract on a slide (report-presentation 5.4).

A deck is not a document either, so PPTX gets a placement: the takeaway becomes
the slide title, the method goes in the slide footer, and the full disclosure —
untruncated — goes in the speaker notes. A slide is the unit that gets pulled out
of a deck and pasted into someone else's, so every slide of a section carries all
three, not just the first.

No exhibit ordinal here, for the same reason as Excel — design decision 7.
"""
import io

from pptx import Presentation

from backend import models, report_builder
from backend.exporters.pptx_exporter import generate_pptx
import pytest

pytestmark = pytest.mark.reporting

_BRANDING = {
    "platform_name": "UKIP",
    "logo_url": None,
    "accent_color": "#6366f1",
    "footer_text": "UKIP",
}


def _seed(db) -> None:
    for idx in range(5):
        db.add(models.RawEntity(
            primary_label=f"Record {idx}", domain="default",
            validation_status="valid" if idx else "pending",
            enrichment_status="completed",
            enrichment_concepts="knowledge graph; ontology",
            enrichment_citation_count=100 + idx,
            enrichment_source="openalex",
            secondary_label="Review" if idx % 2 else "Clinical Trial",
            quality_score=0.8,
        ))
    db.add(models.HarmonizationLog(
        step_id="normalize_labels", step_name="Normalize labels",
        records_updated=5, fields_modified="primary_label",
    ))
    db.add(models.JournalMetric(
        org_id=None, issn_l="issn-x", display_name="Nature Methods",
        normalized_impact_factor=4.10, nif_field="cs", nif_bayes=4.05,
        nif_ci_low=3.60, nif_ci_high=4.55, works_2yr=8, apc_usd=1500,
        is_in_doaj=True,
    ))
    db.commit()


def _deck(db, sections):
    data = generate_pptx(
        db=db, domain_id="default", sections=sections, title=None,
        branding=_BRANDING, org_id=None,
    )
    return Presentation(io.BytesIO(data))


def _slide_text(slide) -> str:
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


def _notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text


def _section_slides(prs, marker: str) -> list:
    """Every slide belonging to one section, found by its label eyebrow."""
    return [s for s in prs.slides if marker in _slide_text(s)]


# ── 5.4 Takeaway, method, notes ───────────────────────────────────────────────


def test_the_slide_is_titled_with_the_finding_not_the_label(db_session):
    _seed(db_session)
    slides = _section_slides(_deck(db_session, ["journal_portfolio"]), "Journal Portfolio")

    assert slides, "journal_portfolio produced no slide"
    text = _slide_text(slides[0])
    assert "in the portfolio" in text, text
    assert "Journal Portfolio" in text, "the label has to remain findable"


def test_the_method_is_on_the_slide_and_in_full_in_the_notes(db_session):
    """The footer is clipped to stay legible; the notes carry the whole thing.

    That split is the requirement: a presenter reads the notes, an audience reads
    the footer, and neither should be the only place the caveat exists.
    """
    _seed(db_session)
    slides = _section_slides(_deck(db_session, ["journal_portfolio"]), "Journal Portfolio")
    # The collector returns a catalog key; the deck renders resolved prose. Compare
    # against what a reader gets, which is what the notes are supposed to carry.
    from backend.reporting.localize import localize_section

    full_method = localize_section(
        report_builder.collect_journal_portfolio(db_session, "default", None), "en"
    ).method

    footer = _slide_text(slides[0])
    assert "field-normalized" in footer, footer

    notes = _notes(slides[0])
    assert full_method in notes, (
        "the notes must carry the disclosure untruncated:\n"
        f"expected: {full_method!r}\ngot: {notes!r}"
    )


def test_the_notes_name_what_the_figure_is_not(db_session):
    """The case the whole disclosure requirement exists for."""
    _seed(db_session)
    slides = _section_slides(_deck(db_session, ["journal_portfolio"]), "Journal Portfolio")
    assert "NOT the Journal Impact Factor" in _notes(slides[0])


def test_every_slide_of_a_section_carries_the_disclosure(db_session):
    """A slide is what gets pulled out of a deck and pasted into another one, so
    the disclosure cannot live only on the first slide of a spilling section."""
    _seed(db_session)
    prs = _deck(db_session, ["authority_control"])
    slides = _section_slides(prs, "Authority Control")

    assert slides
    for idx, slide in enumerate(slides):
        assert _notes(slide), f"slide {idx + 1} of the section has no notes"
        assert "Authority Control" in _slide_text(slide)


def test_slides_do_not_invent_exhibit_ordinals(db_session):
    """Design decision 7 — a deck renders a different set than the document."""
    _seed(db_session)
    prs = _deck(db_session, ["entity_stats", "journal_portfolio"])
    for slide in prs.slides:
        assert "Exhibit" not in _slide_text(slide)
        assert "Exhibit" not in _notes(slide)


# ── The three sections that were bypassing the payload ────────────────────────

#: (section id, slide label, a fragment of its takeaway).
#:
#: `entity_stats` reads "Validation: 4 of 5 …" rather than "4 of 5 entities pass
#: validation": that sentence inflected a noun on one count and a verb on
#: another, which is four whole-sentence variants per language once the copy
#: comes from a catalog. It was rephrased so no word depends on a number.
_MIGRATED = [
    ("entity_stats", "Entity Statistics", "Validation:"),
    ("enrichment_coverage", "Enrichment Coverage", "Enrichment covers"),
    ("top_secondary_labels", "Top Secondary Labels", "leading classification"),
]


def test_the_bespoke_slides_now_go_through_the_payload(db_session):
    """These three had hand-built slides issuing their own queries, so they
    carried no takeaway and no disclosure while the parity map claimed PPTX
    rendered them — the same shape of violation 3.3 found in topic_clusters.

    Migrating cost no detail: the payload is richer than all three hand-built
    slides were (4 KPI cards rather than 2, a Source column, 15 rows rather
    than 10).
    """
    _seed(db_session)
    prs = _deck(db_session, [key for key, _, _ in _MIGRATED])

    for key, label, takeaway_fragment in _MIGRATED:
        slides = _section_slides(prs, label)
        assert slides, f"{key} produced no slide"
        text = _slide_text(slides[0])
        assert takeaway_fragment in text, f"{key} slide states no finding: {text!r}"
        assert _notes(slides[0]), f"{key} slide has no disclosure in its notes"


def test_a_section_is_not_rendered_twice(db_session):
    """The bespoke builders are gone rather than supplemented. If one had been
    left behind, its section would appear on two slides with two titles."""
    _seed(db_session)
    prs = _deck(db_session, ["entity_stats"])

    titled = [s for s in prs.slides if "Entity Statistics" in _slide_text(s)]
    assert len(titled) == 1, f"entity_stats rendered on {len(titled)} slides"
