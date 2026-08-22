"""No rendered report, in any format, may show an unresolved catalog key.

`localize.py` argues that a renderer localizing everything it is handed makes
omission "not expressible". That is true of renderers. It says nothing about
code that reads a payload WITHOUT being one, and #268 found five such readers in
two days:

    * the HTML cover page          (raw f-string, no SectionData)
    * the HTML executive summary   (reads `takeaway` off the collected payload)
    * `_section_manual_note`       (raw HTML, its own default title)
    * Excel's Methodology sheet    (reads takeaway/method off what was collected)
    * the bespoke Harmonization writer (writes `method` into its own cell)

Each was found by a different accident — one by a test asserting on wording that
happened to change, one by CI, three by reading the code. Two of them had been
shipping raw keys in production for weeks.

The per-site fixes do not compose into a guarantee. This does: walk what each
format actually produces and fail on any surface prefix that survived. A section
migrating a field tomorrow is covered without anyone remembering to extend this,
and so is a format added later, because the assertion is about output rather
than about which functions were supposed to call localize_section.

Both a populated and an empty database are exercised. Empty is not redundant:
the executive summary bug only showed on sections with nothing to report, which
is exactly the path a fixture with data never takes.

**What this does NOT prove.** It catches an unresolved KEY, not untranslated
English. Replace a key with a literal English sentence and this stays green —
that is a different defect, and the parity gate plus reading a generated report
are what cover it. A mutation run made the distinction concrete: reverting the
cover's lens caption to the literal "Stakeholder lens" was invisible here, while
reverting the profile label beside it to a raw key was caught immediately.

Mutation-checked against all five known breakages — HTML summary, Excel
Methodology, table cells, the cover's profile label, and the missing comma that
fused two journal keys. Each fails this test when reverted.
"""

from __future__ import annotations

import io
import re

import openpyxl
import pytest
from pptx import Presentation

from backend import models, report_builder
from backend.exporters.excel_exporter import EnterpriseExcelExporter
from backend.exporters.pptx_exporter import generate_pptx
from backend.i18n.catalog import SURFACE_PREFIXES

pytestmark = pytest.mark.reporting

#: Built from the catalog's own prefix list rather than a copy of it, so a new
#: surface is covered the day it is declared.
_KEY_RE = re.compile(
    r"\b(?:%s)[a-z0-9_]+(?:\.[a-z0-9_]+)*"
    % "|".join(re.escape(p) for p in SURFACE_PREFIXES)
)

_SECTIONS = [
    "entity_stats",
    "enrichment_coverage",
    "journal_portfolio",
    "harmonization_log",
    "institutional_benchmark",
    "hidden_patterns",
    "decision_recommendations",
    "impact_projection",
    "authority_control",
    "agentic_trace",
    # Both were missing until #268's fourth batch. Nothing here rendered their
    # copy, so a key leaking from either was invisible to the guard that exists
    # precisely so no format can show one.
    "topic_clusters",
    "collaboration_graph",
]

_BRANDING = {"primary_color": "#2563eb", "organization_name": "UKIP"}


def _seed(db) -> None:
    for idx in range(4):
        db.add(models.RawEntity(
            primary_label=f"Record {idx}", domain="default",
            validation_status="valid" if idx else "pending",
            enrichment_status="completed",
            enrichment_concepts="knowledge graph; ontology",
            enrichment_citation_count=100 + idx,
            enrichment_source="openalex",
            secondary_label="Review",
            quality_score=0.8,
        ))
    db.add(models.HarmonizationLog(
        step_id="normalize_labels", step_name="Normalize labels",
        records_updated=4, fields_modified="primary_label",
    ))
    db.add(models.JournalMetric(
        org_id=None, issn_l="issn-x", display_name="Nature Methods",
        normalized_impact_factor=4.10, nif_field="cs",
        nif_bayes=4.05, nif_ci_low=3.60, nif_ci_high=4.55,
        works_2yr=8, apc_usd=1500, is_in_doaj=True,
    ))
    db.commit()


def _excel_text(data: bytes) -> str:
    wb = openpyxl.load_workbook(io.BytesIO(data))
    return "\n".join(
        str(cell.value)
        for ws in wb.worksheets
        for row in ws.iter_rows()
        for cell in row
        if cell.value is not None
    )


def _pptx_text(data: bytes) -> str:
    deck = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        # Speaker notes carry the untruncated disclosures, so they are output too.
        if slide.has_notes_slide:
            parts.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(parts)


def _every_format(db, language: str) -> dict[str, str]:
    return {
        "html": report_builder.build(
            db, "default", _SECTIONS, org_id=None, language=language
        ),
        "excel": _excel_text(
            EnterpriseExcelExporter().build(db, "default", _SECTIONS, language=language)
        ),
        "pptx": _pptx_text(
            generate_pptx(
                db=db, domain_id="default", sections=_SECTIONS, title=None,
                branding=_BRANDING, org_id=None, language=language,
            )
        ),
    }


@pytest.mark.parametrize("language", ["en", "es"])
@pytest.mark.parametrize("populated", [True, False], ids=["with-data", "empty-db"])
def test_no_format_shows_an_unresolved_catalog_key(db_session, language, populated):
    if populated:
        _seed(db_session)

    for fmt, text in _every_format(db_session, language).items():
        leaked = sorted(set(_KEY_RE.findall(text)))
        assert not leaked, (
            f"{fmt} rendered {len(leaked)} unresolved catalog key(s) in "
            f"{language!r}: {leaked}\n"
            "A payload field became a key and something reads it without going "
            "through localize_section."
        )
