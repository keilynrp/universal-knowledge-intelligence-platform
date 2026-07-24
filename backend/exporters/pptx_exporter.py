"""
PowerPoint exporter — Phase 10 (Artifact Studio)
Generates a branded 16:9 PPTX from report data sections.
Requires python-pptx >= 1.0.2
"""
from __future__ import annotations

import io
from datetime import datetime, timezone
from typing import List, Optional

from sqlalchemy import func
from sqlalchemy.orm import Session

from backend import models
from backend.analyzers.topic_modeling import TopicAnalyzer
from backend.tenant_access import scope_query_to_org

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#rrggbb' to RGBColor."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "6366f1"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def _add_slide(prs: "Presentation") -> any:
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)


def _add_header_bar(slide, accent: "RGBColor", width_emu: int, height_emu: int = 914400 // 8):
    """Add a colored top bar to the slide."""
    bar = slide.shapes.add_shape(1, 0, 0, width_emu, height_emu)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int = 18, bold: bool = False,
                  color: Optional["RGBColor"] = None, wrap: bool = True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def generate_pptx(
    db: Session,
    domain_id: str,
    sections: List[str],
    title: Optional[str],
    branding: dict,
    org_id: int | None = None,
    manual_sections: list[dict[str, str]] | None = None,
) -> bytes:
    """
    Build a branded 16:9 PPTX.
    Returns raw bytes of the .pptx file.
    Raises ImportError if python-pptx is not installed.
    """
    if not _PPTX_AVAILABLE:
        raise ImportError("python-pptx is required for PowerPoint export.")

    # Resolve deprecated aliases (e.g. top_brands) to public ids so every slide
    # gate below matches the vocabulary GET /reports/sections actually returns.
    from backend.report_builder import canonical_sections
    sections = canonical_sections(sections)

    accent = _hex_to_rgb(branding.get("accent_color", "#6366f1"))
    platform = branding.get("platform_name", "UKIP")
    footer_text = branding.get("footer_text", "Universal Knowledge Intelligence Platform")
    report_title = title or f"{platform} Report"

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    entities_query = scope_query_to_org(db.query(models.RawEntity), models.RawEntity, org_id)
    if domain_id:
        entities_query = entities_query.filter(models.RawEntity.domain == domain_id)

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    slide = _add_slide(prs)
    # Full-width accent rectangle (top third)
    cover_bar = slide.shapes.add_shape(1, 0, 0, W, int(H * 0.45))
    cover_bar.fill.solid()
    cover_bar.fill.fore_color.rgb = accent
    cover_bar.line.fill.background()
    # Platform name (white, top area)
    _add_text_box(slide, platform, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
                  font_size=20, bold=True, color=RGBColor(255, 255, 255))
    # Report title (white, larger)
    _add_text_box(slide, report_title, Inches(1), Inches(2.0), Inches(11), Inches(1.1),
                  font_size=32, bold=True, color=RGBColor(255, 255, 255))
    # Date
    date_str = datetime.now(timezone.utc).strftime("%B %d, %Y")
    _add_text_box(slide, date_str, Inches(1), Inches(3.0), Inches(6), Inches(0.4),
                  font_size=13, color=RGBColor(220, 220, 255))
    # Domain
    _add_text_box(slide, f"Domain: {domain_id}", Inches(1), Inches(5.5), Inches(8), Inches(0.4),
                  font_size=12, color=RGBColor(100, 100, 120))
    # Footer
    _add_text_box(slide, footer_text, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
                  font_size=10, color=RGBColor(150, 150, 150))

    for manual in manual_sections or []:
        manual_title = (manual.get("title") or "Analyst Note").strip()[:120]
        manual_content = (manual.get("content") or "").strip()
        if not manual_content:
            continue
        slide = _add_slide(prs)
        _add_header_bar(slide, accent, W)
        _add_text_box(slide, manual_title or "Analyst Note", Inches(0.5), Inches(0.05), Inches(10), Inches(0.45),
                      font_size=16, bold=True, color=RGBColor(255, 255, 255))
        _add_text_box(slide, manual_content[:1800], Inches(0.75), Inches(0.9), Inches(11.8), Inches(5.6),
                      font_size=15, color=RGBColor(45, 55, 72))

    # ── Slide 2: Entity Statistics ────────────────────────────────────────────
    if "entity_stats" in sections:
        total = entities_query.with_entities(func.count(models.RawEntity.id)).scalar() or 0
        by_status = entities_query.with_entities(
            models.RawEntity.validation_status,
            func.count(models.RawEntity.id),
        ).group_by(models.RawEntity.validation_status).all()
        by_enrich = entities_query.with_entities(
            models.RawEntity.enrichment_status,
            func.count(models.RawEntity.id),
        ).group_by(models.RawEntity.enrichment_status).all()
        enrich_map = {r[0]: r[1] for r in by_enrich}
        enriched = enrich_map.get("done", 0) + enrich_map.get("completed", 0)
        enrich_pct = round(enriched / total * 100) if total else 0

        slide = _add_slide(prs)
        _add_header_bar(slide, accent, W)
        _add_text_box(slide, "Entity Statistics", Inches(0.5), Inches(0.05), Inches(10), Inches(0.45),
                      font_size=16, bold=True, color=RGBColor(255, 255, 255))

        stats = [
            ("Total Entities", f"{total:,}"),
            ("Enriched", f"{enriched:,} ({enrich_pct}%)"),
        ]
        for i, (label, value) in enumerate(stats):
            x = Inches(0.5 + i * 4.5)
            box = slide.shapes.add_shape(1, x, Inches(0.8), Inches(3.8), Inches(1.4))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(245, 247, 255)
            box.line.color.rgb = RGBColor(200, 210, 240)
            _add_text_box(slide, label, x + Inches(0.15), Inches(0.85), Inches(3.5), Inches(0.4),
                          font_size=10, color=RGBColor(100, 100, 140))
            _add_text_box(slide, value, x + Inches(0.15), Inches(1.2), Inches(3.5), Inches(0.7),
                          font_size=22, bold=True, color=RGBColor(30, 30, 80))

        # Status breakdown list
        _add_text_box(slide, "Validation Breakdown", Inches(0.5), Inches(2.5), Inches(12), Inches(0.35),
                      font_size=12, bold=True, color=RGBColor(50, 50, 80))
        for idx, (status, count) in enumerate(sorted(by_status, key=lambda x: -x[1])[:6]):
            pct = round(count / total * 100) if total else 0
            line = f"  {status or 'unknown'}:  {count:,}  ({pct}%)"
            _add_text_box(slide, line, Inches(0.5), Inches(2.9 + idx * 0.45), Inches(12), Inches(0.4),
                          font_size=11, color=RGBColor(60, 60, 80))

    # ── Slide 3: Enrichment Coverage ──────────────────────────────────────────
    if "enrichment_coverage" in sections:
        total = entities_query.with_entities(func.count(models.RawEntity.id)).scalar() or 0
        done = entities_query.with_entities(func.count(models.RawEntity.id))\
            .filter(models.RawEntity.enrichment_status.in_(["done", "completed"])).scalar() or 0
        avg_cit = entities_query.with_entities(func.avg(models.RawEntity.enrichment_citation_count))\
            .filter(models.RawEntity.enrichment_status.in_(["done", "completed"])).scalar() or 0
        top_entities = entities_query.with_entities(
            models.RawEntity.primary_label,
            models.RawEntity.enrichment_citation_count,
        )\
            .filter(models.RawEntity.enrichment_status.in_(["done", "completed"]))\
            .order_by(models.RawEntity.enrichment_citation_count.desc()).limit(8).all()
        pct = round(done / total * 100) if total else 0

        slide = _add_slide(prs)
        _add_header_bar(slide, accent, W)
        _add_text_box(slide, "Enrichment Coverage", Inches(0.5), Inches(0.05), Inches(10), Inches(0.45),
                      font_size=16, bold=True, color=RGBColor(255, 255, 255))
        _add_text_box(slide, f"Coverage: {pct}%  |  Enriched: {done:,} of {total:,}  |  Avg citations: {round(avg_cit or 0):,}",
                      Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
                      font_size=13, color=RGBColor(50, 50, 90))
        _add_text_box(slide, "Top Enriched Entities by Citations", Inches(0.5), Inches(1.3), Inches(12), Inches(0.35),
                      font_size=12, bold=True, color=RGBColor(50, 50, 80))
        for idx, (name, cit) in enumerate(top_entities):
            line = f"  {name or '—'}  —  {(cit or 0):,} citations"
            _add_text_box(slide, line, Inches(0.5), Inches(1.75 + idx * 0.45), Inches(12), Inches(0.4),
                          font_size=11, color=RGBColor(60, 60, 80))

    # ── Slide 4: Top Secondary Labels ─────────────────────────────────────────
    if "top_secondary_labels" in sections:
        rows_q = entities_query.with_entities(
            models.RawEntity.secondary_label,
            func.count(models.RawEntity.id).label("n"),
        )\
            .filter(models.RawEntity.secondary_label.isnot(None))\
            .group_by(models.RawEntity.secondary_label)\
            .order_by(func.count(models.RawEntity.id).desc()).limit(10).all()

        slide = _add_slide(prs)
        _add_header_bar(slide, accent, W)
        _add_text_box(slide, "Top Secondary Labels / Classifications", Inches(0.5), Inches(0.05), Inches(10), Inches(0.45),
                      font_size=16, bold=True, color=RGBColor(255, 255, 255))
        total_b = sum(r[1] for r in rows_q)
        for idx, (brand, count) in enumerate(rows_q):
            pct = round(count / total_b * 100) if total_b else 0
            line = f"  {idx+1}.  {brand or '—'}  —  {count:,}  ({pct}%)"
            _add_text_box(slide, line, Inches(0.5), Inches(0.75 + idx * 0.52), Inches(12), Inches(0.45),
                          font_size=12, color=RGBColor(40, 40, 80))

    # ── Slide 5: Topic Clusters ───────────────────────────────────────────────
    if "topic_clusters" in sections:
        try:
            result = TopicAnalyzer().top_topics(domain_id, top_n=20, org_id=org_id)
            topics = result.get("topics", [])
        except Exception:
            topics = []

        slide = _add_slide(prs)
        _add_header_bar(slide, accent, W)
        _add_text_box(slide, "Top Concepts", Inches(0.5), Inches(0.05), Inches(10), Inches(0.45),
                      font_size=16, bold=True, color=RGBColor(255, 255, 255))
        col_size = 10
        left_col  = topics[:col_size]
        right_col = topics[col_size:col_size * 2]
        for idx, t in enumerate(left_col):
            line = f"  {t['concept']}  ({t['count']:,})"
            _add_text_box(slide, line, Inches(0.5), Inches(0.75 + idx * 0.48), Inches(6), Inches(0.42),
                          font_size=11, color=RGBColor(40, 40, 80))
        for idx, t in enumerate(right_col):
            line = f"  {t['concept']}  ({t['count']:,})"
            _add_text_box(slide, line, Inches(6.7), Inches(0.75 + idx * 0.48), Inches(6), Inches(0.42),
                          font_size=11, color=RGBColor(40, 40, 80))

    # ── Final slide: Closing ──────────────────────────────────────────────────
    slide = _add_slide(prs)
    closing_bar = slide.shapes.add_shape(1, 0, 0, W, H)
    closing_bar.fill.solid()
    closing_bar.fill.fore_color.rgb = accent
    closing_bar.line.fill.background()
    _add_text_box(slide, platform, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
                  font_size=28, bold=True, color=RGBColor(255, 255, 255))
    _add_text_box(slide, footer_text, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
                  font_size=14, color=RGBColor(220, 220, 255))
    _add_text_box(slide, f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d')}",
                  Inches(1), Inches(6.8), Inches(6), Inches(0.4),
                  font_size=10, color=RGBColor(200, 200, 230))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
