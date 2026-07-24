"""PPTX renderer over the section payload (unify-report-format-coverage,
phase 2).

Adds one titled slide per section and renders blocks onto it: a StatGrid as a
KPI row of boxes, a Table as a native slide table, a Narrative as a bullet
slide, a Meter as a proportional bar shape.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.presentation import Presentation as PresentationType
from pptx.util import Inches, Pt

from backend.reporting.section_data import (
    Block,
    Meter,
    Narrative,
    SectionData,
    StatGrid,
    Table,
)

SUPPORTED_BLOCKS: frozenset[type] = frozenset({StatGrid, Table, Narrative, Meter})

_INK = RGBColor(0x1E, 0x1E, 0x50)
_MUTED = RGBColor(0x64, 0x64, 0x8C)
_BAR_BG = RGBColor(0xE5, 0xE7, 0xEB)
_CARD_BG = RGBColor(0xF5, 0xF7, 0xFF)


def _blank_slide(prs: PresentationType):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _text(slide, text, left, top, width, height, *, size=12, bold=False, color=_INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = text
    run = frame.paragraphs[0].runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _header(slide, title, accent, width):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    _text(slide, title, Inches(0.5), Inches(0.05), Inches(11), Inches(0.5),
          size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))


def _render_stat_grid(slide, block: StatGrid, top: float) -> None:
    for i, item in enumerate(block.items):
        x = Inches(0.5 + i * 3.0)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(top), Inches(2.8), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = _CARD_BG
        box.line.color.rgb = RGBColor(0xC8, 0xD2, 0xF0)
        _text(slide, item.label, x + Inches(0.15), Inches(top + 0.1), Inches(2.5), Inches(0.35),
              size=10, color=_MUTED)
        _text(slide, item.value, x + Inches(0.15), Inches(top + 0.45), Inches(2.5), Inches(0.6),
              size=20, bold=True)
        if item.sub:
            _text(slide, item.sub, x + Inches(0.15), Inches(top + 1.0), Inches(2.5), Inches(0.3),
                  size=8, color=_MUTED)


def _render_table(slide, block: Table, top: float) -> None:
    rows = len(block.rows) + 1
    cols = len(block.columns)
    shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(top), Inches(12), Inches(0.4 * rows))
    table = shape.table
    for c, name in enumerate(block.columns):
        table.cell(0, c).text = name
    for r, data_row in enumerate(block.rows, start=1):
        for c, cell in enumerate(data_row):
            table.cell(r, c).text = str(cell)


def _render_narrative(slide, block: Narrative, top: float) -> None:
    _text(slide, block.heading, Inches(0.5), Inches(top), Inches(12), Inches(0.5),
          size=15, bold=True)
    body = "\n".join(f"• {p}" for p in block.paragraphs)
    if body:
        _text(slide, body, Inches(0.6), Inches(top + 0.6), Inches(12), Inches(4.5),
              size=13, color=RGBColor(0x2D, 0x37, 0x48))


def _render_meter(slide, block: Meter, top: float) -> None:
    pct = max(0, min(100, round(block.pct)))
    _text(slide, block.label, Inches(0.5), Inches(top), Inches(4), Inches(0.4), size=12, bold=True)
    full = Inches(10)
    track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top + 0.45), full, Inches(0.25))
    track.fill.solid()
    track.fill.fore_color.rgb = _BAR_BG
    track.line.fill.background()
    fill_w = int(full * pct / 100) or 1
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top + 0.45), fill_w, Inches(0.25))
    bar.line.fill.background()
    _text(slide, f"{pct}%", Inches(10.7), Inches(top + 0.4), Inches(1.5), Inches(0.35), size=12, bold=True)


def render_pptx(section: SectionData, prs: PresentationType, accent: RGBColor) -> list:
    """Render one section as slides appended to `prs`. Returns the slides added."""
    slide = _blank_slide(prs)
    added = [slide]
    _header(slide, section.title, accent, prs.slide_width)

    top = 1.0
    for block in section.blocks:
        if top > 6.3:  # spill onto a fresh slide
            slide = _blank_slide(prs)
            added.append(slide)
            _header(slide, section.title, accent, prs.slide_width)
            top = 1.0
        if isinstance(block, StatGrid):
            _render_stat_grid(slide, block, top)
            top += 1.6
        elif isinstance(block, Table):
            _render_table(slide, block, top)
            top += 0.5 + 0.4 * (len(block.rows) + 1)
        elif isinstance(block, Narrative):
            _render_narrative(slide, block, top)
            top += 1.2 + 0.3 * len(block.paragraphs)
        elif isinstance(block, Meter):
            _render_meter(slide, block, top)
            top += 1.0
        else:
            raise TypeError(f"PPTX renderer cannot render {type(block).__name__}")
    return added
